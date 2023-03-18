"""
Script that takes a user text prompt and a .pptx slide deck.
The text prompt must be an instruction explaining how to replace individuals elements in a text sequence,
in this case, a sequence of all the text elements in the .pptx deck."
Calls the OpenAI Edits API to edit the underlying slide xml according to the user's text prompt.
Saves the edited slide deck as a new .pptx file.
"""
import openai
import subprocess
import argparse
import requests
import pptx
import re
import os
from xml.etree import ElementTree as ET
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
from io import BytesIO

def edit_images(save_dir, new_ppt):
    # Edit the images in the slide
    print("   Replacing images...")
    # Open presentation
    prs = pptx.Presentation(new_ppt)
    # Loop through slides
    for i, slide in enumerate(prs.slides):
        prompt = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                prompt += shape.text + " "
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                slide_part, rId = shape.part, shape._element.blip_rId
                image_part = slide_part.related_part(rId)
                # Extract size of old image
                size = image_part.image.size
                if size[0] <= 256:
                    outputsize = "256x256"
                elif size[0] <= 512:
                    outputsize = "512x512"
                else:
                    outputsize = "1024x1024"
                # Call Dall-E API
                print(f"   Calling Dall-E API to generate image for slide {i+1}...")
                response = openai.Image.create(
                    prompt=prompt[:1000],
                    n=1,
                    size=outputsize)
                # Retrieve image from response
                img_url = response["data"][0]["url"]
                img_data = requests.get(img_url).content
                # Resize image and save it
                img = Image.open(BytesIO(img_data))
                img = img.resize(size)
                img.save(os.path.join(save_dir, f"slide{i+1}.png"))
                # Create new image part
                new_img = pptx.parts.image.Image.from_file(os.path.join(save_dir, f"slide{i+1}.png"))
                # overwrite old blob info with new blob info
                image_part.blob = new_img.blob    
    # Save it
    prs.save(os.path.join(save_dir, "final.pptx"))
    print("   ...Images replaced.\n")
                

def edit_text(save_dir, xml_text, response, attrs_per_slide):
        # Edit the text in the slide
        print("   Replacing text...")
        responses = response.split("\n")
        print(f"   {len(responses)} responses.")
        print(f"   {len(xml_text)} text elements to replace.")
        if len(responses) < len(xml_text):
            responses += [""] * (len(xml_text) - len(responses))
        n = 0
        for i, n_slide in enumerate(attrs_per_slide):
            with open(f"{save_dir}/ppt/slides/slide{i+1}.xml", "r+") as outfile:
                content = outfile.read()
                for j in range(n, n+n_slide):
                    content = re.sub(xml_text[j], responses[j], content)
                n += n_slide
                outfile.seek(0)
                outfile.write(content)
                outfile.truncate()
        print("   ...Text edits complete.")

def generate_text(save_dir, num_slides, prompt, temp):
        xml_text = []
        attrs_per_slide = []
        print("   Extracting text from the slide deck...")
        for i in range(1, num_slides+1):
            slide_path = f"{save_dir}/ppt/slides/slide{i}.xml"
            with open(slide_path, "r") as infile:
                xml_string = infile.read()
            slide_text = text_parse(xml_string)
            xml_text += slide_text
            attrs_per_slide.append(len(slide_text))
        inputstr = "\n".join(xml_text)
        print("   Calling the OpenAI Edits API to generate the replacement text...")
        response = openai.Edit.create(
            model="text-davinci-edit-001",
            input=inputstr,
            instruction=prompt + f". The new string must contain {len(xml_text)} lines.",
            temperature=temp)
        #(response["choices"][0]["text"])
        return response["choices"][0]["text"], xml_text, attrs_per_slide

def text_parse(xml_string):
        # Parse the xml string to find the text elements
        tree = ET.fromstring(xml_string)
        a_text = tree.findall(".//a:t", namespaces={"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        elems = [e.text for e in a_text]
        return elems

def count_slides(ppt_file):
        prs = pptx.Presentation(ppt_file)
        num_slides = len(prs.slides)
        return num_slides

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--prompt", help="Text prompt for making edits with OpenAI's Edits API")
    parser.add_argument("--ppt_file", help="Path to the initial slide .pptx file")
    #parser.add_argument("--slide_num", help="Number of the slide in the slide deck to edit")
    parser.add_argument("--save_dir", help="Directory to save the extracted slide deck")
    parser.add_argument("--new_ppt", help="Name of the new slide deck to save")
    #TODO: add temperature argument for API calls
    parser.add_argument("--temperature", help="Temperature for OpenAI's API calls")
    print("\n****** Welcome to the AutoDeck Demo ******\n")
    print("   This demo will take a slide deck and generate new text and images for each slide.")
    print("   The text will be generated using OpenAI's Edits API and the images will be generated using OpenAI's Dall-E 2 API.")
    print("\n...\n")
    # Set the OpenAI API key
    openai.api_key = ""
    # Parse the command line arguments
    args = parser.parse_args()
    # Process the slide deck
    num_slides = count_slides(args.ppt_file)
    # Save the slide deck as .xml
    subprocess.run(["opc", "extract", args.ppt_file, args.save_dir])
    # Generate the new text and images
    response, xml_text, attrs_per_slide = generate_text(args.save_dir, num_slides, args.prompt, float(args.temperature))
    # Edit the text in the slide
    edit_text(args.save_dir, xml_text, response, attrs_per_slide)
    # Repackage the slide deck
    subprocess.run(["opc", "repackage", args.save_dir, args.new_ppt])
    # Edit the images in the slide
    edit_images(args.save_dir, args.new_ppt)
    print("\n****** Deck edits complete! ******\n")
